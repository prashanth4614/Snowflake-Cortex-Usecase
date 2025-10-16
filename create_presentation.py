"""
PowerPoint Generator for Cortex Agent Implementation
Generates a professional PPTX presentation from the project documentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    """Create the complete PowerPoint presentation"""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme (Snowflake theme)
    SNOWFLAKE_BLUE = RGBColor(41, 181, 232)  # #29B5E8
    DARK_BLUE = RGBColor(30, 58, 138)  # #1E3A8A
    GREEN = RGBColor(16, 185, 129)  # #10B981
    DARK_GRAY = RGBColor(31, 41, 55)  # #1F2937
    LIGHT_GRAY = RGBColor(243, 244, 246)  # #F3F4F6
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background color
    background = slide1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = LIGHT_GRAY
    background.line.fill.background()
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Intelligent Sales Assistant"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_BLUE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Powered by Snowflake Cortex AI Agents"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = SNOWFLAKE_BLUE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Key highlights boxes
    highlights = [
        ("🤖 Multi-Tool\nOrchestration", "Automatic coordination\nbetween tools"),
        ("🧵 Conversation\nContext", "Server-side threading\nfor continuity"),
        ("🔄 Auto-Creation", "Self-initializing\nagent"),
        ("⚡ Real-Time\nAnalytics", "SQL generation\nfrom NL")
    ]
    
    box_width = Inches(2)
    box_height = Inches(1.2)
    start_x = Inches(1)
    start_y = Inches(5)
    spacing = Inches(0.3)
    
    for i, (title, desc) in enumerate(highlights):
        x_pos = start_x + (i * (box_width + spacing))
        
        # Box background
        box = slide1.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, start_y, box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = SNOWFLAKE_BLUE
        box.line.width = Pt(2)
        
        # Title
        text_box = slide1.shapes.add_textbox(x_pos, start_y + Inches(0.1), box_width, Inches(0.5))
        tf = text_box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = DARK_BLUE
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Description
        desc_box = slide1.shapes.add_textbox(x_pos, start_y + Inches(0.6), box_width, Inches(0.5))
        df = desc_box.text_frame
        df.text = desc
        df.paragraphs[0].font.size = Pt(9)
        df.paragraphs[0].font.color.rgb = DARK_GRAY
        df.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide 2: Architecture
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    add_slide_title(slide2, "System Architecture", DARK_BLUE)
    
    # Architecture layers
    layers = [
        ("USER INTERFACE LAYER", "Streamlit Application\n• Chat Interface\n• Model Selection\n• Thread Management", Inches(1.5), SNOWFLAKE_BLUE),
        ("ORCHESTRATION LAYER", "CORTEX_SALES_AGENT\n• Query Understanding\n• Tool Selection\n• Response Synthesis", Inches(3.2), DARK_BLUE),
        ("TOOL EXECUTION LAYER", "Cortex Analyst (SQL)\nCortex Search (Docs)", Inches(4.9), GREEN)
    ]
    
    for layer_name, content, y_pos, color in layers:
        # Layer box
        box = slide2.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1), y_pos, Inches(8), Inches(1.2)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = color
        box.line.width = Pt(3)
        
        # Layer title
        title_box = slide2.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.1), Inches(7.6), Inches(0.3))
        tf = title_box.text_frame
        tf.text = layer_name
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = color
        
        # Layer content
        content_box = slide2.shapes.add_textbox(Inches(1.2), y_pos + Inches(0.45), Inches(7.6), Inches(0.7))
        cf = content_box.text_frame
        cf.text = content
        cf.paragraphs[0].font.size = Pt(11)
        cf.paragraphs[0].font.color.rgb = DARK_GRAY
        
        # Add arrow between layers (except last)
        if y_pos < Inches(4.5):
            arrow = slide2.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(4.7), y_pos + Inches(1.3), Inches(0.6), Inches(0.4)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = DARK_GRAY
            arrow.line.fill.background()
    
    # Slide 3: Key Features (Part 1)
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide3, "Key Features & Innovations", DARK_BLUE)
    
    features = [
        {
            "icon": "🤖",
            "title": "Automatic Agent Creation",
            "problem": "Manual setup complexity",
            "solution": "• Auto-detection\n• JSON configuration\n• REST API creation\n• Self-healing",
            "x": Inches(0.5),
            "y": Inches(1.5)
        },
        {
            "icon": "🧵",
            "title": "Intelligent Threading",
            "problem": "Lost conversation context",
            "solution": "• Server-side threads\n• Parent message tracking\n• Pronoun resolution\n• Natural follow-ups",
            "x": Inches(5.2),
            "y": Inches(1.5)
        },
        {
            "icon": "🎯",
            "title": "Smart Query Interpretation",
            "problem": "Ambiguous queries (COUNT vs SELECT)",
            "solution": "• 'list/show' → SELECT\n• 'count' → COUNT(*)\n• 'sum' → SUM()\n• Intent detection",
            "x": Inches(0.5),
            "y": Inches(4.2)
        },
        {
            "icon": "🔄",
            "title": "Multi-Tool Orchestration",
            "problem": "Manual tool selection",
            "solution": "• Automatic routing\n• Parallel execution\n• Response synthesis\n• Unified output",
            "x": Inches(5.2),
            "y": Inches(4.2)
        }
    ]
    
    for feature in features:
        add_feature_box(slide2, feature, SNOWFLAKE_BLUE, DARK_GRAY, GREEN)
    
    # Slide 4: Technical Implementation
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide4, "Technical Implementation", DARK_BLUE)
    
    # Left side - Code flow
    code_box = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5), Inches(4.5), Inches(5)
    )
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(248, 250, 252)
    code_box.line.color.rgb = DARK_BLUE
    code_box.line.width = Pt(2)
    
    code_title = slide4.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(4.1), Inches(0.4))
    ctf = code_title.text_frame
    ctf.text = "Request Flow"
    ctf.paragraphs[0].font.size = Pt(18)
    ctf.paragraphs[0].font.bold = True
    ctf.paragraphs[0].font.color.rgb = DARK_BLUE
    
    code_content = slide4.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(4.1), Inches(4.2))
    ccf = code_content.text_frame
    ccf.text = """User Query
    ↓
snowflake_api_call()
  • Build payload
  • Add thread context
  • POST to agent
    ↓
process_sse_response()
  • Parse events
  • Extract tool results
  • Collect citations
    ↓
Display Results
  • Show response
  • Display SQL
  • Render table"""
    ccf.paragraphs[0].font.size = Pt(11)
    ccf.paragraphs[0].font.name = 'Courier New'
    ccf.paragraphs[0].font.color.rgb = DARK_GRAY
    
    # Right side - Database coverage
    db_box = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.5), Inches(4.3), Inches(5)
    )
    db_box.fill.solid()
    db_box.fill.fore_color.rgb = RGBColor(248, 250, 252)
    db_box.line.color.rgb = GREEN
    db_box.line.width = Pt(2)
    
    db_title = slide4.shapes.add_textbox(Inches(5.4), Inches(1.6), Inches(3.9), Inches(0.4))
    dtf = db_title.text_frame
    dtf.text = "Database Coverage - 9 Tables"
    dtf.paragraphs[0].font.size = Pt(18)
    dtf.paragraphs[0].font.bold = True
    dtf.paragraphs[0].font.color.rgb = GREEN
    
    tables = [
        "CAMPAIGNS - Marketing campaigns",
        "CAMPAIGN_TOUCHES - Interactions",
        "CUSTOMERS - Customer profiles",
        "ORDERS - Sales transactions",
        "ORDER_ITEMS - Line items",
        "PRODUCTS - Product catalog",
        "INVENTORY - Stock levels",
        "REFUNDS - Refund tracking",
        "SHIPMENTS - Delivery tracking"
    ]
    
    db_content = slide4.shapes.add_textbox(Inches(5.4), Inches(2.2), Inches(3.9), Inches(4))
    dcf = db_content.text_frame
    for table in tables:
        p = dcf.add_paragraph()
        p.text = f"• {table}"
        p.font.size = Pt(11)
        p.font.color.rgb = DARK_GRAY
        p.level = 0
    
    # Slide 5: Results & Impact
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide5, "Results & Business Impact", DARK_BLUE)
    
    # Performance metrics
    metrics = [
        ("⏱️ Response Time", "< 3 seconds", GREEN),
        ("🎯 Accuracy", "95%+", GREEN),
        ("🔄 Context", "100%", GREEN),
        ("🛠️ Tools", "Auto", GREEN)
    ]
    
    metric_y = Inches(1.8)
    for i, (label, value, color) in enumerate(metrics):
        x_pos = Inches(0.8) + (i * Inches(2.2))
        
        # Metric box
        box = slide5.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, metric_y, Inches(2), Inches(1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = color
        box.line.width = Pt(3)
        
        # Value
        val_box = slide5.shapes.add_textbox(x_pos, metric_y + Inches(0.15), Inches(2), Inches(0.4))
        vf = val_box.text_frame
        vf.text = value
        vf.paragraphs[0].font.size = Pt(24)
        vf.paragraphs[0].font.bold = True
        vf.paragraphs[0].font.color.rgb = color
        vf.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Label
        lab_box = slide5.shapes.add_textbox(x_pos, metric_y + Inches(0.6), Inches(2), Inches(0.3))
        lf = lab_box.text_frame
        lf.text = label
        lf.paragraphs[0].font.size = Pt(11)
        lf.paragraphs[0].font.color.rgb = DARK_GRAY
        lf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Before/After comparison
    comparison = slide5.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(3.2), Inches(8.4), Inches(2.8)
    )
    comparison.fill.solid()
    comparison.fill.fore_color.rgb = LIGHT_GRAY
    comparison.line.fill.background()
    
    # Table header
    headers = ["Aspect", "Before", "After", "Improvement"]
    header_y = Inches(3.3)
    col_widths = [Inches(2.5), Inches(2), Inches(2), Inches(1.9)]
    x_start = Inches(0.9)
    
    for i, header in enumerate(headers):
        x_pos = x_start + sum(col_widths[:i])
        hbox = slide5.shapes.add_textbox(x_pos, header_y, col_widths[i], Inches(0.3))
        hf = hbox.text_frame
        hf.text = header
        hf.paragraphs[0].font.size = Pt(12)
        hf.paragraphs[0].font.bold = True
        hf.paragraphs[0].font.color.rgb = DARK_BLUE
    
    # Table rows
    rows = [
        ("Query Method", "Manual SQL", "Natural language", "100%"),
        ("Tool Selection", "Manual", "Automatic", "100%"),
        ("Context", "Restart each", "Continuous", "100%"),
        ("Setup Time", "Hours", "Minutes", "90%")
    ]
    
    row_y_start = Inches(3.7)
    for i, row in enumerate(rows):
        row_y = row_y_start + (i * Inches(0.45))
        for j, cell in enumerate(row):
            x_pos = x_start + sum(col_widths[:j])
            cbox = slide5.shapes.add_textbox(x_pos, row_y, col_widths[j], Inches(0.4))
            cf = cbox.text_frame
            cf.text = cell
            cf.paragraphs[0].font.size = Pt(10)
            cf.paragraphs[0].font.color.rgb = DARK_GRAY
            if j == 3:  # Improvement column
                cf.paragraphs[0].font.color.rgb = GREEN
                cf.paragraphs[0].font.bold = True
    
    # Slide 6: ROI & Conclusion
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide6, "ROI & Strategic Value", DARK_BLUE)
    
    # ROI Section
    roi_box = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.8), Inches(4), Inches(2.5)
    )
    roi_box.fill.solid()
    roi_box.fill.fore_color.rgb = RGBColor(240, 253, 244)
    roi_box.line.color.rgb = GREEN
    roi_box.line.width = Pt(3)
    
    roi_title = slide6.shapes.add_textbox(Inches(1), Inches(1.9), Inches(3.6), Inches(0.4))
    rtf = roi_title.text_frame
    rtf.text = "Time Savings"
    rtf.paragraphs[0].font.size = Pt(18)
    rtf.paragraphs[0].font.bold = True
    rtf.paragraphs[0].font.color.rgb = GREEN
    
    roi_content = slide6.shapes.add_textbox(Inches(1), Inches(2.4), Inches(3.6), Inches(1.8))
    rcf = roi_content.text_frame
    rcf.text = "⏱️ 80% reduction in time\n   to get sales insights\n\n⏱️ 90% reduction in\n   policy lookup time\n\n⏱️ 70% reduction in\n   training time"
    rcf.paragraphs[0].font.size = Pt(13)
    rcf.paragraphs[0].font.color.rgb = DARK_GRAY
    
    # Success Factors
    success_box = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.8), Inches(4), Inches(2.5)
    )
    success_box.fill.solid()
    success_box.fill.fore_color.rgb = RGBColor(239, 246, 255)
    success_box.line.color.rgb = SNOWFLAKE_BLUE
    success_box.line.width = Pt(3)
    
    success_title = slide6.shapes.add_textbox(Inches(5.4), Inches(1.9), Inches(3.6), Inches(0.4))
    stf = success_title.text_frame
    stf.text = "Success Factors"
    stf.paragraphs[0].font.size = Pt(18)
    stf.paragraphs[0].font.bold = True
    stf.paragraphs[0].font.color.rgb = SNOWFLAKE_BLUE
    
    success_content = slide6.shapes.add_textbox(Inches(5.4), Inches(2.4), Inches(3.6), Inches(1.8))
    scf = success_content.text_frame
    scf.text = "✅ Native Snowflake Features\n\n✅ User-Centric Design\n\n✅ Robust Engineering\n\n✅ Enterprise Standards"
    scf.paragraphs[0].font.size = Pt(13)
    scf.paragraphs[0].font.color.rgb = DARK_GRAY
    
    # Strategic Advantages
    strat_box = slide6.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(4.6), Inches(8.4), Inches(1.8)
    )
    strat_box.fill.solid()
    strat_box.fill.fore_color.rgb = LIGHT_GRAY
    strat_box.line.fill.background()
    
    strat_title = slide6.shapes.add_textbox(Inches(1), Inches(4.7), Inches(8), Inches(0.4))
    sttf = strat_title.text_frame
    sttf.text = "Strategic Advantages"
    sttf.paragraphs[0].font.size = Pt(16)
    sttf.paragraphs[0].font.bold = True
    sttf.paragraphs[0].font.color.rgb = DARK_BLUE
    
    strat_content = slide6.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(1.1))
    stcf = strat_content.text_frame
    stcf.text = "🎯 Democratized data access across organization\n🎯 Consistent data interpretation and reporting\n🎯 Scalable foundation for AI-driven insights"
    stcf.paragraphs[0].font.size = Pt(13)
    stcf.paragraphs[0].font.color.rgb = DARK_GRAY
    
    # Save presentation
    prs.save('Cortex_Agent_Presentation.pptx')
    print("✅ Presentation created successfully: Cortex_Agent_Presentation.pptx")

def add_slide_title(slide, title_text, color):
    """Add a title to a slide"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = color

def add_feature_box(slide, feature, primary_color, text_color, accent_color):
    """Add a feature box to a slide"""
    # Main box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        feature["x"], feature["y"], Inches(4.5), Inches(2.3)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    box.line.color.rgb = primary_color
    box.line.width = Pt(2)
    
    # Icon
    icon_box = slide.shapes.add_textbox(feature["x"] + Inches(0.2), feature["y"] + Inches(0.1), Inches(0.5), Inches(0.5))
    icon_frame = icon_box.text_frame
    icon_frame.text = feature["icon"]
    icon_frame.paragraphs[0].font.size = Pt(32)
    
    # Title
    title_box = slide.shapes.add_textbox(feature["x"] + Inches(0.8), feature["y"] + Inches(0.15), Inches(3.5), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = feature["title"]
    title_frame.paragraphs[0].font.size = Pt(14)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = primary_color
    
    # Problem
    prob_box = slide.shapes.add_textbox(feature["x"] + Inches(0.2), feature["y"] + Inches(0.7), Inches(4.1), Inches(0.3))
    prob_frame = prob_box.text_frame
    prob_frame.text = f"Problem: {feature['problem']}"
    prob_frame.paragraphs[0].font.size = Pt(10)
    prob_frame.paragraphs[0].font.italic = True
    prob_frame.paragraphs[0].font.color.rgb = text_color
    
    # Solution
    sol_box = slide.shapes.add_textbox(feature["x"] + Inches(0.2), feature["y"] + Inches(1.1), Inches(4.1), Inches(1.1))
    sol_frame = sol_box.text_frame
    sol_frame.text = f"Solution:\n{feature['solution']}"
    sol_frame.paragraphs[0].font.size = Pt(10)
    sol_frame.paragraphs[0].font.color.rgb = text_color

if __name__ == "__main__":
    try:
        create_presentation()
    except Exception as e:
        print(f"❌ Error creating presentation: {str(e)}")
        import traceback
        traceback.print_exc()
